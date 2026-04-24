import os
import tempfile
from typing import List, Dict, Any

import streamlit as st
import pandas as pd
import ollama
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LCDocument

# ------------------------------------------------------------
# Check for optional embedding packages
# ------------------------------------------------------------
try:
    from langchain_huggingface import HuggingFaceEmbeddings
    from langchain_community.vectorstores import FAISS
    from langchain_community.retrievers import BM25Retriever
    EMBEDDINGS_AVAILABLE = True
except ImportError:
    EMBEDDINGS_AVAILABLE = False

# ------------------------------------------------------------
# Streamlit UI setup
# ------------------------------------------------------------
st.set_page_config(page_title="RAG Assistant for PDF Summarization", layout="wide")
st.title(" PDF Assistant")
st.markdown("Upload documents, ask questions, get answers using **Mistral** + **Hybrid Search** (FAISS + BM25).")

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB

# ------------------------------------------------------------
# Session state initialisation
# ------------------------------------------------------------
if "doc_store" not in st.session_state:
    st.session_state.doc_store = {}          # filename -> {"faiss": FAISS, "bm25": BM25Retriever}
if "memory_store" not in st.session_state:
    st.session_state.memory_store = {}       # filename or "ALL" -> list of Q&A pairs
if "messages" not in st.session_state:
    st.session_state.messages = []

# ------------------------------------------------------------
# Helper functions
# ------------------------------------------------------------
def load_file_from_bytes(file_bytes, filename: str) -> str:
    """Extract text from uploaded file bytes."""
    ext = filename.split('.')[-1].lower()
    try:
        if ext == 'pdf':
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            reader = PdfReader(tmp_path)
            text = "".join(p.extract_text() or "" for p in reader.pages)
            os.unlink(tmp_path)
            return text
        elif ext == 'txt':
            return file_bytes.decode("utf-8")
        elif ext == 'csv':
            import io
            df = pd.read_csv(io.BytesIO(file_bytes))
            return df.to_string()
        elif ext == 'docx':
            import io
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs)
        elif ext == 'pptx':
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return "\n".join(text)
        else:
            return ""
    except Exception as e:
        st.error(f"File reading error: {e}")
        return ""

def split_text(text: str, chunk_size=500, chunk_overlap=50) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)

def build_index(chunks: List[str]) -> Dict[str, Any]:
    """Build FAISS (dense) and BM25 (sparse) indexes."""
    if not EMBEDDINGS_AVAILABLE:
        st.error("⚠️ Missing embedding packages. Run:\n`pip install sentence-transformers faiss-cpu langchain-community langchain-huggingface rank-bm25`")
        return None
    if not chunks:
        st.error("No text chunks to index.")
        return None

    with st.spinner("Building search index (this may take a moment)..."):
        # Convert to LangChain documents
        lc_docs = [LCDocument(page_content=chunk) for chunk in chunks]
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        faiss_index = FAISS.from_documents(lc_docs, embeddings)
        bm25_retriever = BM25Retriever.from_documents(lc_docs)
        bm25_retriever.k = 3
    return {"faiss": faiss_index, "bm25": bm25_retriever, "chunks_count": len(chunks)}

def hybrid_search(query: str, indexes: List[Dict], top_k=5) -> List[str]:
    """Combine results from FAISS and BM25 using .invoke() for BM25."""
    all_docs = []
    for idx in indexes:
        # FAISS similarity search
        faiss_results = idx["faiss"].similarity_search(query, k=3)
        all_docs.extend(faiss_results)

        # BM25: use .invoke() (new API) or fallback to get_relevant_documents
        bm25 = idx["bm25"]
        try:
            if hasattr(bm25, "invoke"):
                bm25_results = bm25.invoke(query)
            elif hasattr(bm25, "get_relevant_documents"):
                bm25_results = bm25.get_relevant_documents(query)
            else:
                bm25_results = []
        except Exception:
            bm25_results = []
        all_docs.extend(bm25_results)

    # Deduplicate by content
    seen = set()
    unique = []
    for doc in all_docs:
        content = doc.page_content.strip()
        if content and content not in seen:
            unique.append(content)
            seen.add(content)
    return unique[:top_k]

def build_prompt(query: str, context: str, memory: List[str]) -> str:
    memory_text = "\n".join(memory[-3:]) if memory else ""
    return f"""
Recent conversation:
{memory_text}

Relevant context from document(s):
{context if context else "No relevant context found."}

Question: {query}

Answer (be concise and helpful):
"""

def answer_stream(query: str, doc_name: str = None):
    """Generator that streams the answer."""
    if not EMBEDDINGS_AVAILABLE:
        yield "⚠️ Embeddings not available. Please install the required packages."
        return

    # Determine which indexes to use
    if doc_name and doc_name in st.session_state.doc_store:
        indexes = [st.session_state.doc_store[doc_name]]
        memory_key = doc_name
    else:
        indexes = list(st.session_state.doc_store.values())
        memory_key = "ALL"

    if not indexes:
        yield "❌ No documents uploaded. Please upload a file first."
        return

    # Retrieve context
    context_chunks = hybrid_search(query, indexes)
    context = "\n\n---\n\n".join(context_chunks)

    # Get memory
    memory = st.session_state.memory_store.get(memory_key, [])

    # Build prompt
    prompt = build_prompt(query, context, memory)

    # Stream from Ollama
    try:
        stream = ollama.chat(
            model="mistral",
            messages=[{"role": "user", "content": prompt}],
            stream=True
        )
        full_answer = ""
        for chunk in stream:
            text_chunk = chunk["message"]["content"]
            full_answer += text_chunk
            yield text_chunk

        # Update memory
        st.session_state.memory_store.setdefault(memory_key, []).append(f"Q: {query}\nA: {full_answer}")

    except Exception as e:
        yield f"\n⚠️ Ollama error: {str(e)}\n\nMake sure Ollama is running (`ollama serve`) and the mistral model is pulled (`ollama pull mistral`)."

# ------------------------------------------------------------
# Sidebar – File upload and document selection
# ------------------------------------------------------------
with st.sidebar:
    st.header("📂 Upload Documents")

    uploaded_file = st.file_uploader(
        "Choose a file (max 50 MB)",
        type=['pdf', 'txt', 'csv', 'docx', 'pptx']
    )

    if uploaded_file is not None:
        if uploaded_file.size > MAX_FILE_SIZE:
            st.error(f"File too large! Maximum size is {MAX_FILE_SIZE // (1024*1024)} MB.")
        else:
            with st.spinner(f"Processing {uploaded_file.name}..."):
                file_bytes = uploaded_file.read()
                text = load_file_from_bytes(file_bytes, uploaded_file.name)
                if text:
                    chunks = split_text(text)
                    index = build_index(chunks)
                    if index:
                        st.session_state.doc_store[uploaded_file.name] = index
                        st.success(f"✅ {uploaded_file.name} added ({index['chunks_count']} chunks)")
                    else:
                        st.error("Failed to build index.")
                else:
                    st.error("No text could be extracted.")

    st.divider()
    doc_list = list(st.session_state.doc_store.keys())
    if doc_list:
        selected_doc = st.selectbox("Select document to query", ["ALL"] + doc_list)
    else:
        selected_doc = "ALL"
        st.info("No documents uploaded yet.")

    if st.button("🗑️ Clear all documents"):
        st.session_state.doc_store = {}
        st.session_state.memory_store = {}
        st.session_state.messages = []
        st.rerun()

# ------------------------------------------------------------
# Main chat area
# ------------------------------------------------------------
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

if prompt := st.chat_input("Ask a question about your document(s)"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        placeholder = st.empty()
        full_response = ""
        doc_param = None if selected_doc == "ALL" else selected_doc
        for chunk in answer_stream(prompt, doc_param):
            full_response += chunk
            placeholder.markdown(full_response + "▌")
        placeholder.markdown(full_response)

    st.session_state.messages.append({"role": "assistant", "content": full_response})
